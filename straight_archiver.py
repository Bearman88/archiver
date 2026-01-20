import os
import re
import shutil
import csv
from datetime import datetime

from docx import Document
import pdfplumber
from pdf2image import convert_from_path
import pytesseract
from pptx import Presentation

# ================= 配置区 =================
ROOT_DIR = r"E:\sata11-17703836269\国网不错的报告" # 改这个
OCR_LANG = "chi_sim+eng"
MAX_TITLE_LEN = 60

DRY_RUN = False
ENABLE_ROLLBACK = True

MOVE_LOG = "move_log.csv"
ERROR_LOG = "error.log"

# ---------------- 扩展名 ----------------
IMAGE_EXT = [".png", ".jpg", ".jpeg"]
EXCEL_EXT = [".xlsx", ".xls"]
OTHER_DOC_EXT = [".doc", ".wps"]
DATA_EXT = [".txt", ".csv", ".dta"]
SCRIPT_EXT = [".py", ".do", ".r", ".m"]
GRAPH_EXT = [".gph"]
ARCHIVE_EXT = [".zip", ".rar", ".7z"]
CAJ_EXT = [".caj"]
EXCLUDE_EXT = [".exe", ".bat", ".msi"]

# ---------------- 行政文种规则 ----------------
TOPIC_RULES = [
    ("会议纪要", ["会议纪要", "会议记录", "纪要"]),
    ("请示", ["请示"]),
    ("决定", ["决定", "决议", "批复"]),
    ("通知", ["通知", "通告", "公告","函"]),
    ("报告", ["报告", "研究", "汇报" ,"说明", "材料","分析"]),
]

# ---------------- 参考资料子分类 ----------------
REFERENCE_SUBCATEGORIES = {
    "论文": ["论文", "paper", "article", "journal", "thesis", "academic", "research article", "scientific"],
    "报告": ["报告", "report", "research report", "technical report", "汇报", "summary"],
    "合同": ["合同", "agreement", "contract", "契约", "deal", "purchase agreement"],
    "标准": ["标准", "规范", "specification", "guideline", "protocol", "准则", "standard"]
}

ILLEGAL = r'[\\/:*?"<>|]'

# ================= 日志 =================
def log_error(msg):
    ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    with open(ERROR_LOG, "a", encoding="utf-8") as f:
        f.write(f"[{ts}] {msg}\n")

def log_move(old, new, is_folder=False):
    if DRY_RUN or not ENABLE_ROLLBACK:
        return
    exists = os.path.exists(MOVE_LOG)
    with open(MOVE_LOG, "a", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        if not exists:
            w.writerow(["time", "old_path", "new_path", "is_folder"])
        w.writerow([datetime.now().isoformat(), old, new, int(is_folder)])

# ================= 工具 =================
def clean(text):
    if not text:
        return ""
    text = re.sub(ILLEGAL, "", text.strip())
    text = re.sub(r"\s+", "", text)
    return text[:MAX_TITLE_LEN]

def get_date(path):
    t = min(os.path.getctime(path), os.path.getmtime(path))
    return datetime.fromtimestamp(t).strftime("%Y%m%d")

def resolve_dup(path):
    base, ext = os.path.splitext(path)
    i = 1
    new = path
    while os.path.exists(new):
        new = f"{base}({i}){ext}"
        i += 1
    return new

# ================= 标题提取 =================
def pick_title(lines, fallback_name=""):
    """优先中文最大行，再英文"""
    for l in lines:
        if 5 <= len(l) <= MAX_TITLE_LEN and re.search(r'[\u4e00-\u9fff]', l):
            return clean(l)
    for l in lines:
        if 5 <= len(l) <= MAX_TITLE_LEN and re.search(r'[A-Za-z]', l):
            return clean(l)
    if fallback_name and re.search(r'[\u4e00-\u9fff]', fallback_name):
        return clean(fallback_name)
    if lines:
        return clean(lines[0])
    return "_未识别标题_"

# ================= 文档提取 =================
def extract_docx(path):
    try:
        doc = Document(path)
        lines = [p.text for p in doc.paragraphs if p.text.strip()]
        return pick_title(lines, fallback_name=os.path.splitext(os.path.basename(path))[0])
    except Exception as e:
        log_error(f"DOCX失败: {path} | {e}")
        return pick_title([], fallback_name=os.path.splitext(os.path.basename(path))[0])

def extract_pdf(path):
    lines = []
    try:
        with pdfplumber.open(path) as pdf:
            page = pdf.pages[0]
            # 按top分行
            rows = {}
            for c in page.chars:
                key = round(c['top'], 1)
                rows.setdefault(key, []).append(c)
            best_text, best_size = "", 0
            for chars in rows.values():
                text = "".join(c['text'] for c in chars).strip()
                if not text or len(text) < 5:
                    continue
                avg_size = sum(c['size'] for c in chars)/len(chars)
                if avg_size > best_size:
                    best_size = avg_size
                    best_text = text
            if best_text:
                return clean(best_text)
            # fallback 全页文本
            text = page.extract_text() or ""
            lines.extend([l.strip() for l in text.splitlines() if l.strip()])
    except Exception as e:
        log_error(f"PDF文本失败: {path} | {e}")

    # OCR兜底
    if not lines:
        try:
            imgs = convert_from_path(path, first_page=1, last_page=1)
            txt = pytesseract.image_to_string(imgs[0], lang=OCR_LANG)
            lines.extend([l.strip() for l in txt.splitlines() if l.strip()])
        except Exception as e:
            log_error(f"PDF OCR失败: {path} | {e}")
    return pick_title(lines, fallback_name=os.path.splitext(os.path.basename(path))[0])

def extract_caj(path):
    lines = []
    try:
        imgs = convert_from_path(path, first_page=1, last_page=1)
        for img in imgs:
            txt = pytesseract.image_to_string(img, lang=OCR_LANG)
            lines.extend([l.strip() for l in txt.splitlines() if l.strip()])
    except Exception as e:
        log_error(f"CAJ提取失败: {path} | {e}")
    return pick_title(lines, fallback_name=os.path.splitext(os.path.basename(path))[0])

def extract_pptx(path):
    try:
        prs = Presentation(path)
        slide = prs.slides[0]
        best_text, best_size = "", 0
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size and r.font.size.pt > best_size:
                        text = p.text.strip()
                        if 5 <= len(text) <= MAX_TITLE_LEN:
                            best_size = r.font.size.pt
                            best_text = text
        return pick_title([best_text] if best_text else [], fallback_name=os.path.splitext(os.path.basename(path))[0])
    except Exception as e:
        log_error(f"PPTX失败: {path} | {e}")
        return pick_title([], fallback_name=os.path.splitext(os.path.basename(path))[0])

# ================= 参考资料子类 =================
def detect_reference_subcategory(path, filename, ext):
    base = filename.lower()
    for subcat, kws in REFERENCE_SUBCATEGORIES.items():
        for k in kws:
            if k.lower() in base:
                return f"参考资料/{subcat}"
    return "参考资料/其他参考资料"

# ================= 文种识别 =================
def detect_topic(title, filename, ext):
    if ext in [".doc", ".docx", ".wps"]:
        for topic, kws in TOPIC_RULES:
            for k in kws:
                if k.lower() in (title+filename).lower():
                    return topic
        return "其他"
    elif ext in CAJ_EXT + [".pdf"]:
        return detect_reference_subcategory(os.path.join(ROOT_DIR, filename), filename, ext)
    else:
        return None

# ================= 删除空文件夹 =================
def remove_empty_dirs(path):
    for root, dirs, files in os.walk(path, topdown=False):
        for d in dirs:
            dir_path = os.path.join(root, d)
            if not os.listdir(dir_path):
                os.rmdir(dir_path)
                print(f"[DEL] 空文件夹删除: {dir_path}")
                log_move(dir_path, "__FOLDER_DELETED__", is_folder=True)

# ================= 主流程 =================
def organize():
    print("=== 文件整理开始 ===")
    for root, dirs, files in os.walk(ROOT_DIR):
        for f in files:
            old = os.path.join(root, f)
            name, ext = os.path.splitext(f)
            ext = ext.lower()

            title, category = "", None

            if ext == ".docx":
                title = extract_docx(old)
                category = detect_topic(title, f, ext)
            elif ext == ".pdf":
                title = extract_pdf(old)
                category = detect_reference_subcategory(old, f, ext)
            elif ext == ".pptx":
                title = extract_pptx(old)
                category = "PPT文件"
            elif ext in CAJ_EXT:
                title = extract_caj(old)
                category = detect_reference_subcategory(old, f, ext)
            elif ext in IMAGE_EXT:
                category = "图像文件"
            elif ext in EXCEL_EXT:
                category = "Excel文件"
            elif ext in OTHER_DOC_EXT:
                title = extract_docx(old)
                category = detect_topic(title, f, ext)
            elif ext in DATA_EXT:
                category = "数据集文件"
            elif ext in SCRIPT_EXT:
                category = "脚本程序"
            elif ext in GRAPH_EXT:
                category = "图形文件"
            elif ext in ARCHIVE_EXT:
                category = "压缩包"
            elif ext in EXCLUDE_EXT:
                category = "人工判别"
            else:
                category = "其他"

            date = get_date(old)
            title_for_name = title or os.path.splitext(f)[0]
            new_name = f"{date}_{title_for_name}{ext}"

            target_dir = os.path.join(ROOT_DIR, category)
            os.makedirs(target_dir, exist_ok=True)
            new_path = resolve_dup(os.path.join(target_dir, new_name))

            if DRY_RUN:
                print(f"[DRY] {f} → {category}/{os.path.basename(new_path)}")
            else:
                shutil.move(old, new_path)
                log_move(old, new_path)
                print(f"[OK]  {f} → {category}/{os.path.basename(new_path)}")

    remove_empty_dirs(ROOT_DIR)
    print("=== 文件整理结束 ===")

# ================= 回滚 =================
def rollback():
    if not os.path.exists(MOVE_LOG):
        print("无回滚日志")
        return
    print("=== 回滚开始 ===")
    with open(MOVE_LOG, newline="", encoding="utf-8") as f:
        rows = list(csv.DictReader(f))
    for r in reversed(rows):
        old = r["old_path"]
        new = r["new_path"]
        is_folder = r.get("is_folder") == "1"
        if is_folder:
            if not os.path.exists(old):
                os.makedirs(old)
                print(f"[ROLLBACK] 恢复空文件夹: {old}")
        else:
            if not os.path.exists(new):
                print(f"[WARN] 文件不存在，跳过: {new}")
                continue
            restore = resolve_dup(old)
            os.makedirs(os.path.dirname(restore), exist_ok=True)
            shutil.move(new, restore)
            print(f"[ROLLBACK] {new} → {restore}")
    print("=== 回滚完成 ===")

# ================= 入口 =================
if __name__ == "__main__":
    organize()
    # rollback()
