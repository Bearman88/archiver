import os
from docx import Document
import pdfplumber
from pptx import Presentation

def extract_text(path):
    ext = os.path.splitext(path)[1].lower()

    try:
        if ext == '.docx':
            doc = Document(path)
            return [p.text.strip() for p in doc.paragraphs if p.text.strip()][:5]

        if ext == '.pdf':
            with pdfplumber.open(path) as pdf:
                text = pdf.pages[0].extract_text()
                return text.splitlines()[:5] if text else []

        if ext == '.pptx':
            prs = Presentation(path)
            slide = prs.slides[0]
            lines = []
            for s in slide.shapes:
                if s.has_text_frame:
                    lines.append(s.text.strip())
            return lines

    except Exception:
        return []

    return []